import fitz
import os
import subprocess
from pptx import Presentation
from docx import Document

def convert_ppt_to_pptx(file_path: str) -> str:
    output_dir = "/tmp"
    subprocess.run([
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "--headless",
        "--convert-to", "pptx",
        "--outdir", output_dir,
        file_path
    ])
    new_path = os.path.join(output_dir, os.path.basename(file_path).replace(".ppt", ".pptx"))
    return new_path

def extract_text(file_path: str) -> str:
    if file_path.endswith(".pdf"):
        return extract_pdf(file_path)
    elif file_path.endswith(".pptx"):
        return extract_pptx(file_path)
    elif file_path.endswith(".ppt"):
        pptx_path = convert_ppt_to_pptx(file_path)
        return extract_pptx(pptx_path)
    elif file_path.endswith(".docx"):
        return extract_docx(file_path)
    elif file_path.endswith(".txt"):
        with open(file_path, "r") as f:
            return f.read()
    else:
        return ""

def extract_pdf(file_path: str) -> str:
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    return text

def extract_pptx(file_path: str) -> str:
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_docx(file_path: str) -> str:
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text